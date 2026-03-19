import os
import sys
import json
import base64
import time
from pathlib import Path
import chromadb
from core import embeddings as _emb
import pymupdf4llm
import fitz  # PyMuPDF — bundled with pymupdf4llm
from docx import Document
import pandas as pd
from bs4 import BeautifulSoup
import markdown
from datetime import datetime

# ---------------------------------------------------------------------------
# Vision-enhanced ingestion prompt
#
# Sent to the vision model once per PDF page.  The model reads the rendered
# page image and produces a structured plain-text description that is stored
# as the primary content of the chunk.  This description drives embedding
# quality; the raw extracted text is appended for exact keyword fallback.
# ---------------------------------------------------------------------------
import re as _re  # used by _strip_thinking below

# ---------------------------------------------------------------------------
# Metadata extraction prompt
#
# Sent to the LLM after each chunk (vision-enhanced or plain text) when
# --mdextraction is passed.  The output JSON is appended to the chunk as
# [METADATA: {...}] and key fields are stored in ChromaDB metadata for
# structured lookup and multi-hop retrieval.
# ---------------------------------------------------------------------------
METADATA_EXTRACT_PROMPT = (
    "You are indexing a reference document page. Extract all named entities as structured JSON.\n\n"
    "Target entities: creatures, characters, spells, powers, abilities, items, rules, stats, locations.\n"
    "For each entity extract: name, type, and any attributes/powers/stats visible on this page.\n\n"
    "Return ONLY valid JSON in this exact format, nothing else:\n"
    '{"entities": [{"name": "...", "type": "...", "attributes": {}, "powers": []}]}\n\n'
    'If no named entities are present, return: {"entities": []}\n'
    "Do NOT include any explanation, markdown, or text outside the JSON."
)


def _strip_thinking(text: str) -> str:
    """Remove <thinking>...</thinking> blocks from LLM output before parsing."""
    return _re.sub(r'<thinking>.*?</thinking>', '', text, flags=_re.DOTALL).strip()


VISION_EXTRACT_PROMPT = (
    "You are indexing a page from a reference document. Extract all content precisely.\n\n"
    "Rules:\n"
    "- Stat blocks / tables: list EVERY entry with ALL values — "
    "name, category, and every column value exactly as shown.\n"
    "- Rules / mechanics text: state the rule clearly and list every named game element.\n"
    "- Prose / narrative: one concise sentence summary only.\n"
    "- Blank / decorative pages: respond with just the word SKIP.\n\n"
    "Be complete and exact. Include every name and number visible on the page."
)


class DocumentIngester:
    """Ingests documents into ChromaDB for RAG"""

    def __init__(self, documents_dir=None, chroma_path=None):
        # Try to load paths from config
        config_path = Path("config/settings.json")
        embed_model = "BAAI/bge-base-en-v1.5"
        if config_path.exists():
            with open(config_path) as f:
                config = json.load(f)
            documents_dir = documents_dir or config.get("documents", {}).get("directory", "documents")
            chroma_path = chroma_path or config.get("memory", {}).get("chroma_path", "data/chroma_db")
            embed_model = config.get("memory", {}).get("embedding_model", embed_model)
        else:
            documents_dir = documents_dir or "documents"
            chroma_path = chroma_path or "data/chroma_db"

        self._embed_model = embed_model

        self.documents_dir = Path(documents_dir)
        self.chroma_path = chroma_path

        self.client = chromadb.PersistentClient(path=chroma_path)
        self.collection = self.client.get_or_create_collection(
            name="talon_documents",
            metadata={"description": "User documents for RAG retrieval"}
        )

        self.supported_extensions = {
            '.pdf', '.txt', '.docx', '.pptx', '.md', '.markdown',
            '.py', '.js', '.java', '.cpp', '.cs', '.html', '.htm',
            '.csv', '.xlsx', '.xls', '.epub'
        }

        print(f"✓ Document ingester ready!")
        print(f"  Documents directory: {self.documents_dir}")
        print(f"  Supported types: {', '.join(self.supported_extensions)}\n")

    def chunk_text(self, text, chunk_size=400, overlap=50):
        """Split text into overlapping chunks.

        Default chunk_size raised to 400 words (from 200) to keep cross-
        referenced content (e.g. Mana Bolt described under Manaball) within
        the same chunk, improving RAG retrieval on structured reference books.
        """
        words = text.split()
        chunks = []

        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if chunk:
                chunks.append(chunk)

        return chunks

    def extract_pdf(self, filepath):
        """Extract text from PDF.

        Primary: pymupdf4llm.to_markdown() — preserves reading order in
        multi-column layouts, tables, and headings.

        Fallback: raw fitz page-by-page extraction — used when pymupdf4llm
        returns suspiciously little content (some PDFs with complex layouts
        cause pymupdf4llm to silently drop most pages).
        """
        import fitz  # PyMuPDF — bundled with pymupdf4llm

        # Primary: pymupdf4llm
        try:
            result = pymupdf4llm.to_markdown(str(filepath))
            doc = fitz.open(str(filepath))
            page_count = doc.page_count
            doc.close()
            # Sanity check: ~100 chars per page minimum
            if result and len(result.strip()) >= page_count * 100:
                return result
            print(f"    ⚠ pymupdf4llm returned only {len(result or '')} chars "
                  f"for {page_count} pages — falling back to raw fitz extraction")
        except Exception as e:
            print(f"    ⚠ pymupdf4llm failed ({e}) — falling back to raw fitz")

        # Fallback: raw fitz page-by-page
        try:
            doc = fitz.open(str(filepath))
            pages = []
            for page in doc:
                text = page.get_text(flags=fitz.TEXT_PRESERVE_WHITESPACE)
                if text.strip():
                    pages.append(text)
            doc.close()
            full_text = "\n\n".join(pages)
            print(f"    → fitz fallback: extracted {len(full_text)} chars "
                  f"from {len(pages)} pages")
            return full_text if full_text.strip() else None
        except Exception as e:
            print(f"  ✗ fitz fallback also failed: {e}")
            return None

    def extract_docx(self, filepath):
        """Extract text from DOCX"""
        try:
            doc = Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except Exception as e:
            print(f"  ✗ Error reading DOCX: {e}")
            return None

    def extract_pptx(self, filepath):
        """Extract text from PowerPoint (.pptx) — one block per slide."""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
            return "\n\n".join(parts) if parts else None
        except Exception as e:
            print(f"  ✗ Error reading PPTX: {e}")
            return None

    def extract_txt(self, filepath):
        """Extract text from TXT/code files"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            print(f"  ✗ Error reading text file: {e}")
            return None

    def extract_markdown(self, filepath):
        """Extract text from Markdown"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                md_text = file.read()
                html = markdown.markdown(md_text)
                soup = BeautifulSoup(html, 'html.parser')
                return soup.get_text()
        except Exception as e:
            print(f"  ✗ Error reading Markdown: {e}")
            return None

    def extract_csv_excel(self, filepath):
        """Extract data from CSV/Excel"""
        try:
            if filepath.suffix == '.csv':
                df = pd.read_csv(filepath)
            else:
                df = pd.read_excel(filepath)

            text = f"File: {filepath.name}\n\n"
            text += f"Columns: {', '.join(df.columns)}\n\n"
            text += df.to_string()
            return text
        except Exception as e:
            print(f"  ✗ Error reading data file: {e}")
            return None

    def extract_html(self, filepath):
        """Extract text from HTML"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text()
        except Exception as e:
            print(f"  ✗ Error reading HTML: {e}")
            return None

    @staticmethod
    def _html_chapter_to_text(html_bytes: bytes) -> str:
        """Convert an EPUB chapter's HTML to structured plain text.

        Preserves table structure as markdown so technical content (stats,
        spec tables, comparison grids) survives ingestion.

        Two-pass approach: tables are converted to markdown and replaced
        in-place BEFORE any iteration, avoiding BeautifulSoup linked-list
        corruption from decompose() during a descendants walk.
        """
        from bs4 import NavigableString as _NS
        soup = BeautifulSoup(html_bytes, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()

        # Pass 1: replace every <table> with its markdown equivalent in-place.
        # find_all() returns a static list so replacements don't affect iteration.
        for table in soup.find_all("table"):
            rows = []
            for tr in table.find_all("tr"):
                cells = [td.get_text(" ", strip=True)
                         for td in tr.find_all(["td", "th"])]
                if cells:
                    rows.append("| " + " | ".join(cells) + " |")
            if rows:
                sep = "| " + " | ".join(
                    ["---"] * max(rows[0].count("|") - 1, 1)) + " |"
                rows.insert(1, sep)
                table.replace_with(_NS("\n" + "\n".join(rows) + "\n"))
            else:
                table.decompose()

        # Pass 2: plain text extraction — tables are now inline markdown strings.
        root = soup.body if soup.body else soup
        text = root.get_text(separator="\n")
        return "\n".join(line for line in text.splitlines() if line.strip())

    def extract_epub(self, filepath):
        """Extract text from EPUB with structured table preservation.

        Requires: pip install ebooklib beautifulsoup4
        """
        try:
            import ebooklib
            from ebooklib import epub as _epub
        except ImportError:
            print("  ✗ ebooklib not installed — run: pip install ebooklib")
            return None

        try:
            book = _epub.read_epub(str(filepath), options={"ignore_ncx": True})
            parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                text = self._html_chapter_to_text(item.get_content())
                if text:
                    parts.append(text)
            return "\n\n".join(parts) if parts else None
        except Exception as e:
            print(f"  ✗ Error reading EPUB: {e}")
            return None

    def ingest_epub_with_vision(self, filepath, llm_client,
                                chunk_size: int = 400, overlap: int = 50,
                                md_extract: bool = False) -> int:
        """Vision-enhanced EPUB ingestion.

        For each chapter: stores structured text (tables preserved as markdown).
        For each embedded image: describes with vision LLM and stores as a
        separate chunk, mirroring the PDF vision pipeline.

        Chunk format for image chunks:
            [VISION: <model description>]

            SOURCE: <epub filename>, embedded image <n>
        """
        try:
            import ebooklib
            from ebooklib import epub as _epub
        except ImportError:
            print("  ✗ ebooklib not installed — run: pip install ebooklib")
            return 0

        print(f"  Processing (EPUB+vision): {filepath.name}")

        try:
            book = _epub.read_epub(str(filepath), options={"ignore_ncx": True})
        except Exception as exc:
            print(f"  ✗ Could not open EPUB: {exc}")
            return 0

        # Clear stale chunks
        try:
            existing = self.collection.get(where={"filename": filepath.name}, include=[])
            if existing["ids"]:
                self.collection.delete(ids=existing["ids"])
                print(f"    → Cleared {len(existing['ids'])} existing chunks")
        except Exception as exc:
            print(f"    ⚠ Could not clear existing chunks: {exc}")

        doc_id_base = f"doc_{filepath.stem}_{int(datetime.now().timestamp())}"
        chunks_stored = 0
        img_index = 0
        t_start = time.time()

        # ── Text chapters — chunked like the text pipeline ────────────────────
        chapters = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))
        for ch_idx, item in enumerate(chapters):
            try:
                chapter_text = self._html_chapter_to_text(item.get_content())
            except Exception as exc:
                print(f"    ⚠ Chapter {ch_idx} extraction failed: {exc}")
                continue
            if not chapter_text or len(chapter_text.split()) < 15:
                continue

            # Split long chapters into overlapping word-count chunks so
            # retrieval works at paragraph granularity, not chapter granularity.
            words = chapter_text.split()
            sub_idx = 0
            for start in range(0, len(words), chunk_size - overlap):
                sub_text = " ".join(words[start:start + chunk_size])
                if len(sub_text.split()) < 15:
                    continue
                meta = {
                    "filename": filepath.name,
                    "source": str(filepath),
                    "chapter": ch_idx,
                    "sub_chunk": sub_idx,
                    "type": "epub_chapter",
                }
                doc_id = f"{doc_id_base}_ch{ch_idx}_s{sub_idx}"
                self.collection.add(
                    embeddings=_emb.embed_documents([sub_text], self._embed_model),
                    documents=[sub_text],
                    metadatas=[meta],
                    ids=[doc_id],
                )
                chunks_stored += 1
                sub_idx += 1

        # ── Embedded images — raster only, skip SVG/tiny ──────────────────────
        _RASTER_TYPES = {"image/jpeg", "image/png", "image/gif", "image/webp"}
        all_images = list(book.get_items_of_type(ebooklib.ITEM_IMAGE))
        raster_images = [
            img for img in all_images
            if img.media_type in _RASTER_TYPES and len(img.get_content()) >= 1024
        ]
        print(f"    → {len(chapters)} chapters, "
              f"{len(raster_images)}/{len(all_images)} raster image(s) to describe")

        for img_item in raster_images:
            img_b64 = base64.b64encode(img_item.get_content()).decode()
            img_index += 1
            print(f"    [img {img_index}/{len(raster_images)}] Describing...",
                  end=" ", flush=True)

            t_img = time.time()
            vision_desc = self._describe_page(llm_client, img_b64,
                                              img_index, filepath.name)
            elapsed_img = time.time() - t_img
            print(f"{elapsed_img:.1f}s")

            if not vision_desc:
                continue

            chunk_text = (
                f"[VISION: {vision_desc}]\n\n"
                f"SOURCE: {filepath.name}, embedded image {img_index}"
            )
            meta = {
                "filename": filepath.name,
                "source": str(filepath),
                "image_index": img_index,
                "type": "epub_image",
            }
            doc_id = f"{doc_id_base}_img{img_index}"
            self.collection.add(
                embeddings=_emb.embed_documents([chunk_text], self._embed_model),
                documents=[chunk_text],
                metadatas=[meta],
                ids=[doc_id],
            )
            chunks_stored += 1

        elapsed_total = time.time() - t_start
        print(f"    ✓ Ingested {chunks_stored} chunks "
              f"in {int(elapsed_total // 60)}m{int(elapsed_total % 60):02d}s "
              f"(epub+vision)")
        return chunks_stored

    def extract_text(self, filepath):
        """Route to appropriate extractor"""
        ext = filepath.suffix.lower()

        if ext == '.pdf':
            return self.extract_pdf(filepath)
        elif ext == '.docx':
            return self.extract_docx(filepath)
        elif ext == '.pptx':
            return self.extract_pptx(filepath)
        elif ext in ['.txt', '.py', '.js', '.java', '.cpp', '.cs']:
            return self.extract_txt(filepath)
        elif ext in ['.md', '.markdown']:
            return self.extract_markdown(filepath)
        elif ext in ['.csv', '.xlsx', '.xls']:
            return self.extract_csv_excel(filepath)
        elif ext in ['.html', '.htm']:
            return self.extract_html(filepath)
        elif ext == '.epub':
            return self.extract_epub(filepath)
        else:
            return None

    # ------------------------------------------------------------------
    # Vision-enhanced PDF ingestion
    # ------------------------------------------------------------------

    def _extract_metadata(self, llm_client, chunk_text: str,
                          page_num: int, filename: str) -> dict | None:
        """Call the LLM to extract structured entity metadata from chunk text.

        Args:
            llm_client:  Running LLMClient instance.
            chunk_text:  Chunk text to analyse (vision-enhanced or plain).
            page_num:    1-based page/chunk number for log messages.
            filename:    Filename for log messages.

        Returns:
            Parsed dict with 'entities' list, or None on any failure.
        """
        try:
            prompt = METADATA_EXTRACT_PROMPT + "\n\nText to index:\n" + chunk_text[:1500]
            raw = llm_client.generate(prompt, max_length=512, temperature=0.0)
            raw = _strip_thinking(raw)
            # Extract the outermost JSON object from the response
            start = raw.index("{")
            end   = raw.rindex("}") + 1
            parsed = json.loads(raw[start:end])
            if isinstance(parsed, dict) and isinstance(parsed.get("entities"), list):
                return parsed
            return None
        except Exception as exc:
            print(f"      ⚠ Metadata extraction failed p{page_num} {filename}: {exc}")
            return None

    def _describe_page(self, llm_client, img_b64: str, page_num: int,
                       filename: str) -> str:
        """Call the vision model to describe a single rendered PDF page.

        Returns a plain-text description string, or "" on failure.
        The model is instructed to return "SKIP" for blank/decorative pages
        so we can drop those without storing empty chunks.
        """
        try:
            desc = llm_client.generate(
                VISION_EXTRACT_PROMPT,
                use_vision=True,
                screenshot_b64=img_b64,
                temperature=0.0,
            )
            if not desc or desc.startswith("Error:"):
                print(f"      ⚠ Vision failed p{page_num}: {desc}")
                return ""
            if desc.strip().upper() == "SKIP":
                return ""
            return desc.strip()
        except Exception as exc:
            print(f"      ⚠ Vision exception p{page_num}: {exc}")
            return ""

    def ingest_pdf_with_vision(self, filepath, llm_client,
                               chunk_size: int = 400, overlap: int = 50,
                               md_extract: bool = False) -> int:
        """Vision-enhanced PDF ingestion — one chunk per page.

        Each chunk is structured as:

            [VISION: <model description of the page>]

            RAW TEXT:
            <fitz plain-text extraction>

        The VISION section produces a rich, semantically meaningful embedding
        (spell names + exact stats, rule names + descriptions, etc.).
        The RAW TEXT section preserves exact strings for keyword/$contains
        fallback searches.

        Very long pages (>800 words combined) are sub-split: the vision header
        is repeated on every sub-chunk so embedding quality is maintained.

        Args:
            filepath:    Path to the PDF file.
            llm_client:  Instantiated LLMClient pointing at the running server.
            chunk_size:  Word count target for raw-text sub-splits (if needed).
            overlap:     Overlap words for raw-text sub-splits.

        Returns:
            Number of chunks stored in ChromaDB.
        """
        print(f"  Processing (vision): {filepath.name}")

        doc = fitz.open(str(filepath))
        page_count = doc.page_count
        print(f"    → {page_count} pages to describe")

        # Remove stale chunks for this file before re-ingesting
        try:
            existing = self.collection.get(
                where={"filename": filepath.name},
                include=[],
            )
            if existing["ids"]:
                self.collection.delete(ids=existing["ids"])
                print(f"    → Cleared {len(existing['ids'])} existing chunks")
        except Exception as exc:
            print(f"    ⚠ Could not clear existing chunks: {exc}")

        doc_id_base = f"doc_{filepath.stem}_{int(datetime.now().timestamp())}"
        chunks_stored = 0
        t_start = time.time()

        for page_idx in range(page_count):
            page = doc[page_idx]
            page_num = page_idx + 1  # 1-based for display

            # --- Raw text for this page (keyword fallback) ---
            raw_text = page.get_text().strip()

            # --- ETA display ---
            elapsed = time.time() - t_start
            if page_idx > 0:
                rate = elapsed / page_idx          # seconds per page so far
                remaining = rate * (page_count - page_idx)
                eta_str = f"  ETA ~{int(remaining // 60)}m{int(remaining % 60):02d}s"
            else:
                eta_str = ""
            print(f"    [{page_num:>4}/{page_count}] Describing...{eta_str}",
                  end=" ", flush=True)

            # Skip vision for image-heavy / art pages — fewer than 15 words of
            # extractable text means it's almost certainly a full-page artwork,
            # chapter divider, or blank page.  These render as huge pixmaps and
            # flood KoboldCpp with image tokens causing stalls or timeouts.
            if len(raw_text.split()) < 15:
                print(f"0.0s  (art/blank — vision skipped)")
                vision_desc = ""
            else:
                # Render page to base64 PNG only when we know it has content
                pixmap = page.get_pixmap(dpi=100)
                img_b64 = base64.b64encode(pixmap.tobytes("png")).decode()
                t_page = time.time()
                vision_desc = self._describe_page(llm_client, img_b64, page_num,
                                                  filepath.name)
                page_elapsed = time.time() - t_page
                print(f"{page_elapsed:.1f}s")

            # Skip genuinely empty pages
            if not vision_desc and not raw_text:
                print(f"      ↳ skipped (blank page)")
                continue

            if vision_desc:
                print(f"      ↳ {vision_desc[:100].replace(chr(10), ' ')}")

            # --- Build enriched chunk text ---
            parts = []
            if vision_desc:
                parts.append(f"[VISION: {vision_desc}]")
            if raw_text:
                parts.append(f"RAW TEXT:\n{raw_text}")
            chunk_text = "\n\n".join(parts)

            # --- Structured entity metadata extraction (--mdextraction) ---
            base_meta = {
                "type": "document",
                "filename": filepath.name,
                "filepath": str(filepath),
                "file_type": filepath.suffix,
                "page_number": page_idx,
                "total_pages": page_count,
                "vision_enhanced": True,
                "ingested_at": datetime.now().isoformat(),
            }
            if md_extract:
                entity_meta = self._extract_metadata(
                    llm_client, chunk_text, page_num, filepath.name)
                if entity_meta and entity_meta.get("entities"):
                    entities = entity_meta["entities"]
                    base_meta["entity_names"] = ", ".join(
                        e["name"] for e in entities if e.get("name"))
                    base_meta["entity_types"] = ", ".join(
                        e["type"] for e in entities if e.get("type"))
                    base_meta["has_metadata"] = True
                    chunk_text = chunk_text + (
                        f"\n\n[METADATA: {json.dumps(entity_meta, ensure_ascii=False)}]"
                    )

            # --- Sub-split only if the combined page is very long ---

            if len(chunk_text.split()) > 800 and raw_text:
                # Repeat the vision header on every sub-chunk so each one
                # retains full semantic context in its embedding.
                vision_header = (f"[VISION: {vision_desc}]\n\n"
                                 if vision_desc else "")
                sub_chunks = self.chunk_text(raw_text,
                                             chunk_size=chunk_size,
                                             overlap=overlap)
                for sub_i, raw_sub in enumerate(sub_chunks):
                    sub_text = f"{vision_header}RAW TEXT:\n{raw_sub}"
                    doc_id = f"{doc_id_base}_p{page_idx}_s{sub_i}"
                    meta = {**base_meta, "sub_chunk": sub_i}
                    self.collection.add(
                        embeddings=_emb.embed_documents([sub_text], self._embed_model),
                        documents=[sub_text],
                        metadatas=[meta],
                        ids=[doc_id],
                    )
                    chunks_stored += 1
            else:
                doc_id = f"{doc_id_base}_p{page_idx}"
                meta = {**base_meta, "sub_chunk": 0}
                self.collection.add(
                    embeddings=_emb.embed_documents([chunk_text], self._embed_model),
                    documents=[chunk_text],
                    metadatas=[meta],
                    ids=[doc_id],
                )
                chunks_stored += 1

        doc.close()
        total_elapsed = time.time() - t_start
        print(f"    ✓ Ingested {chunks_stored} page-chunks "
              f"in {int(total_elapsed // 60)}m{int(total_elapsed % 60):02d}s "
              f"(vision-enhanced)")
        return chunks_stored

    def ingest_file(self, filepath, chunk_size=400, overlap=50,
                    llm_client=None, use_vision=False, md_extract=False):
        """Ingest a single file.

        Args:
            filepath:    Path to the file.
            chunk_size:  Words per chunk (text path).
            overlap:     Overlap words between chunks (text path).
            llm_client:  LLMClient instance (required for --vision/--mdextraction).
            use_vision:  If True and file is a PDF, use vision-enhanced ingestion.
            md_extract:  If True, run LLM entity metadata extraction per chunk.
        """
        # Vision path: PDF or EPUB + use_vision flag + running LLM server
        if use_vision and llm_client is not None:
            if filepath.suffix.lower() == ".pdf":
                return self.ingest_pdf_with_vision(
                    filepath, llm_client, chunk_size=chunk_size, overlap=overlap,
                    md_extract=md_extract)
            if filepath.suffix.lower() == ".epub":
                return self.ingest_epub_with_vision(
                    filepath, llm_client, chunk_size=chunk_size, overlap=overlap,
                    md_extract=md_extract)

        print(f"  Processing: {filepath.name}")

        text = self.extract_text(filepath)
        if not text or len(text.strip()) < 50:
            print(f"    ⚠ Skipped (no content or too short)")
            return 0

        # Remove any existing chunks for this file so re-ingest is idempotent
        # and doesn't accumulate duplicate chunks across multiple runs.
        try:
            existing = self.collection.get(
                where={"filename": filepath.name},
                include=[],
            )
            if existing["ids"]:
                self.collection.delete(ids=existing["ids"])
                print(f"    → Replaced {len(existing['ids'])} existing chunks")
        except Exception as e:
            print(f"    ⚠ Could not delete existing chunks: {e}")

        chunks = self.chunk_text(text, chunk_size=chunk_size, overlap=overlap)
        print(f"    → Created {len(chunks)} chunks ({chunk_size}w/{overlap}w overlap)")

        doc_id_base = f"doc_{filepath.stem}_{int(datetime.now().timestamp())}"

        for i, chunk in enumerate(chunks):
            doc_id = f"{doc_id_base}_chunk_{i}"
            chunk_text = chunk

            meta: dict = {
                "type": "document",
                "filename": filepath.name,
                "filepath": str(filepath),
                "file_type": filepath.suffix,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "ingested_at": datetime.now().isoformat(),
            }

            # Structured entity metadata extraction (--mdextraction, text path)
            if md_extract and llm_client is not None:
                entity_meta = self._extract_metadata(
                    llm_client, chunk_text, i + 1, filepath.name)
                if entity_meta and entity_meta.get("entities"):
                    entities = entity_meta["entities"]
                    meta["entity_names"] = ", ".join(
                        e["name"] for e in entities if e.get("name"))
                    meta["entity_types"] = ", ".join(
                        e["type"] for e in entities if e.get("type"))
                    meta["has_metadata"] = True
                    chunk_text = chunk_text + (
                        f"\n\n[METADATA: {json.dumps(entity_meta, ensure_ascii=False)}]"
                    )

            self.collection.add(
                embeddings=_emb.embed_documents([chunk_text], self._embed_model),
                documents=[chunk_text],
                metadatas=[meta],
                ids=[doc_id],
            )

        print(f"    ✓ Ingested {len(chunks)} chunks")
        return len(chunks)

    def ingest_directory(self, chunk_size=400, overlap=50, llm_client=None,
                         use_vision=False, md_extract=False):
        """Ingest all supported files.

        Args:
            chunk_size:  Words per chunk (text path).
            overlap:     Overlap words between chunks (text path).
            llm_client:  LLMClient instance (required for --vision/--mdextraction).
            use_vision:  If True, PDFs use the vision-enhanced page-by-page path.
            md_extract:  If True, run LLM entity metadata extraction per chunk.
        """
        if not self.documents_dir.exists():
            print(f"Creating documents directory: {self.documents_dir}")
            self.documents_dir.mkdir(parents=True, exist_ok=True)
            print(f"\n✓ Created! Add your documents to: {self.documents_dir}")
            return

        files = []
        for ext in self.supported_extensions:
            files.extend(self.documents_dir.rglob(f"*{ext}"))

        if not files:
            print(f"No documents found in {self.documents_dir}")
            return

        mode_parts = []
        if use_vision:
            mode_parts.append("vision")
        if md_extract:
            mode_parts.append("mdextraction")
        mode_note = f" ({', '.join(mode_parts)})" if mode_parts else ""
        print(f"Found {len(files)} files to ingest{mode_note}\n")

        total_chunks = 0
        successful = 0

        for filepath in files:
            try:
                chunks = self.ingest_file(
                    filepath,
                    chunk_size=chunk_size,
                    overlap=overlap,
                    llm_client=llm_client,
                    use_vision=use_vision,
                    md_extract=md_extract,
                )
                total_chunks += chunks
                successful += 1
            except Exception as e:
                print(f"  ✗ Failed: {e}")

        print(f"\n{'=' * 50}")
        print(f"✓ Ingestion complete!")
        print(f"  Files: {successful}/{len(files)}")
        print(f"  Chunks: {total_chunks}")
        print(f"{'=' * 50}\n")

    def list_documents(self):
        """List all ingested documents"""
        results = self.collection.get()

        if not results['ids']:
            print("No documents in database")
            return

        files = {}
        for metadata in results['metadatas']:
            filename = metadata.get('filename', 'unknown')
            if filename not in files:
                files[filename] = {
                    'type': metadata.get('file_type', 'unknown'),
                    'chunks': 0,
                    'ingested_at': metadata.get('ingested_at', 'unknown')
                }
            files[filename]['chunks'] += 1

        print(f"\nDocuments in database: {len(files)}\n")

        for filename, info in sorted(files.items()):
            print(f"  {filename}")
            print(f"    Type: {info['type']}, Chunks: {info['chunks']}")

    def clear_documents(self):
        """Clear all documents"""
        self.client.delete_collection(name="talon_documents")
        self.collection = self.client.get_or_create_collection(
            name="talon_documents",
            metadata={"description": "User documents for RAG retrieval"}
        )
        print("✓ Cleared all documents")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Talon document ingester",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python ingest_documents.py                              # plain text ingest\n"
            "  python ingest_documents.py --vision                     # vision-enhanced PDFs\n"
            "  python ingest_documents.py --mdextraction               # entity metadata extraction\n"
            "  python ingest_documents.py --vision --mdextraction      # both\n"
            "  python ingest_documents.py list                         # list indexed documents\n"
            "  python ingest_documents.py clear                        # wipe the database\n"
        ),
    )
    parser.add_argument("command", nargs="?", default="ingest",
                        choices=["ingest", "list", "clear"],
                        help="Command to run (default: ingest)")
    parser.add_argument("--chunk-size", type=int, default=400,
                        help="Words per chunk for text path (default: 400)")
    parser.add_argument("--overlap", type=int, default=50,
                        help="Overlap words between chunks (default: 50)")
    parser.add_argument("--vision", action="store_true",
                        help=(
                            "Use vision model to describe each PDF page before "
                            "chunking. Produces richer embeddings for structured "
                            "content (tables, stat blocks). Requires LLM server "
                            "with mmproj to be running."
                        ))
    parser.add_argument("--mdextraction", action="store_true",
                        help=(
                            "Extract structured entity metadata (names, types, "
                            "attributes) from each chunk via LLM and store as "
                            "[METADATA: {json}]. Enables multi-hop retrieval. "
                            "Can be combined with --vision. Requires LLM server."
                        ))
    args = parser.parse_args()

    ingester = DocumentIngester()

    if args.command == "list":
        ingester.list_documents()
    elif args.command == "clear":
        confirm = input("Clear all documents? (yes/no): ")
        if confirm.lower() == "yes":
            ingester.clear_documents()
    else:
        llm_client = None

        if args.vision or args.mdextraction:
            # ----------------------------------------------------------------
            # Vision / mdextraction mode: spin up LLMClient from settings.json
            # so we can call the running LLM server during ingestion.
            # ----------------------------------------------------------------
            sys.path.insert(0, str(Path(__file__).parent))
            from core.llm_client import LLMClient  # noqa: E402

            cfg_path = Path("config/settings.json")
            if cfg_path.exists():
                with open(cfg_path) as f:
                    cfg = json.load(f)
            else:
                # Sensible defaults if settings.json is missing
                cfg = {
                    "llm": {
                        "endpoint": "http://localhost:5001/api/v1/generate",
                        "max_length": 512,
                        "temperature": 0.1,
                        "top_p": 0.9,
                        "rep_pen": 1.1,
                        "timeout": 120,
                        "stop_sequences": ["<|im_end|>", "<|im_start|>"],
                        "prompt_template": {
                            "user_prefix": "<|im_start|>user\n",
                            "user_suffix": "<|im_end|>\n",
                            "assistant_prefix": "<|im_start|>assistant\n",
                            "vision_prefix": (
                                "<|vision_start|><|image_pad|><|vision_end|>"
                            ),
                        },
                        "api_format": "koboldcpp",
                    }
                }

            mode_label = "+".join(
                m for m, f in [("vision", args.vision), ("mdextraction", args.mdextraction)] if f
            )
            print(f"{mode_label} mode — connecting to LLM server...")
            llm_client = LLMClient(cfg)
            if not llm_client.test_connection():
                print(
                    "⚠  Could not reach LLM server.  "
                    "Falling back to plain text extraction.\n"
                    "   Ensure your LLM server is running before using "
                    "--vision or --mdextraction.\n"
                )
                llm_client = None
            else:
                print(f"✓ LLM server ready — {mode_label} ingestion enabled\n")

        mode_parts = []
        if args.vision and llm_client:
            mode_parts.append("vision")
        if args.mdextraction and llm_client:
            mode_parts.append("mdextraction")
        mode = "+".join(mode_parts) if mode_parts else "text"
        print(f"Mode: {mode}  |  chunk-size: {args.chunk_size}  |  overlap: {args.overlap}\n")
        ingester.ingest_directory(
            chunk_size=args.chunk_size,
            overlap=args.overlap,
            llm_client=llm_client,
            use_vision=bool(args.vision and llm_client),
            md_extract=bool(args.mdextraction and llm_client),
        )