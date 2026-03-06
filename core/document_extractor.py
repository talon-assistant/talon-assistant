"""document_extractor.py — Extract plain text from common document formats for LLM review.

Used when a user attaches a non-image file to the GUI chat.  The extracted text is
injected into the LLM prompt so the model can summarise, answer questions, etc.

Supported formats (all optional-import, graceful if library missing):
  .pdf              PyMuPDF (fitz)        — already required by ingest_documents.py
  .docx             python-docx           — already in requirements.txt
  .xlsx / .xls      pandas + openpyxl     — pandas already in requirements.txt
  .pptx             python-pptx           — requirements.txt
  .txt / .md / .rst built-in open()
  .csv              pandas                — already in requirements.txt

This is for one-shot inline review, NOT ChromaDB ingestion — see
ingest_documents.py for the batch RAG pipeline.
"""

import os

MAX_CHARS = 32_000  # ~8 000 tokens — leaves ~4 000 tokens for overhead + response
                    # in a 12 288-token context. Raise if you have a larger context.

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif", ".tiff"}
_DOC_EXTS   = {".pdf", ".docx", ".xlsx", ".xls", ".pptx",
               ".txt", ".md", ".rst", ".csv"}


def is_image(path: str) -> bool:
    """Return True if the file extension indicates a raster image."""
    return os.path.splitext(path)[1].lower() in _IMAGE_EXTS


def is_document(path: str) -> bool:
    """Return True if the file extension is a supported document format."""
    return os.path.splitext(path)[1].lower() in _DOC_EXTS


def extract(path: str) -> str | None:
    """Extract plain text from a document file.

    Returns the extracted text (trimmed to MAX_CHARS) or None on failure.
    The returned string does NOT include structural markers — callers should
    wrap it with _wrap_external() before injecting into an LLM prompt.
    """
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        if ext == ".docx":
            return _extract_docx(path)
        if ext in (".xlsx", ".xls"):
            return _extract_excel(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext == ".csv":
            return _extract_csv(path)
        if ext in (".txt", ".md", ".rst"):
            return _extract_text(path)
        return None
    except Exception as e:
        print(f"   [DocExtractor] Failed to extract '{os.path.basename(path)}': {e}")
        return None


# ── Private helpers ────────────────────────────────────────────────────────────

def _trim(text: str) -> str:
    text = text.strip()
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n\n[... document truncated at {MAX_CHARS:,} chars — attach a smaller section or ingest via ingest_documents.py for full RAG access ...]"
    return text


def _extract_pdf(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    parts = []
    for i, page in enumerate(doc):
        t = page.get_text().strip()
        if t:
            parts.append(f"[Page {i + 1}]\n{t}")
    doc.close()
    return _trim("\n\n".join(parts))


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    # Include table content
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return _trim("\n\n".join(parts))


def _extract_excel(path: str) -> str:
    import pandas as pd
    xl = pd.ExcelFile(path)
    parts = []
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        parts.append(f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}")
    return _trim("\n\n".join(parts))


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
    return _trim("\n\n".join(parts))


def _extract_csv(path: str) -> str:
    import pandas as pd
    df = pd.read_csv(path)
    return _trim(df.to_string(index=False))


def _extract_text(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return _trim(f.read())
